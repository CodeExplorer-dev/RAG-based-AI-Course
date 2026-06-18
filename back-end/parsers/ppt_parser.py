"""PPT/PPTX 课件解析器
使用 python-pptx 逐页提取文本，识别幻灯片标题作为章节标题，
返回带标题信息的页面列表。
"""

from pptx import Presentation


class PPTParser:
    """PPT 解析：遍历每个 slide，识别标题与正文内容"""

    def extract(self, file_path: str) -> list[dict]:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            title_text = ''
            body_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                full_text = shape.text_frame.text.strip()
                if not full_text:
                    continue
                # 通常 PPT 模板的第一个文本框是标题
                if shape == slide.shapes[0] and shape.has_text_frame:
                    title_text = full_text
                else:
                    body_texts.append(full_text)
            # 合并
            content = ''
            if title_text:
                content = f'【{title_text}】\n'
            if body_texts:
                content += '\n'.join(body_texts)
            if content:
                slides.append({
                    'page_number': i,
                    'text': content,
                    'heading': title_text or None,
                })
        return slides
