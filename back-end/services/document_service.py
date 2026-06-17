"""文档处理服务 — 文件解析、文本提取、分块调度"""

import os
import re
import logging
from typing import List
from werkzeug.utils import secure_filename
from werkzeug.datastructures import FileStorage
from extensions import db
from models import Courseware, DocumentChunk
from utils.errors import NotFoundError

logger = logging.getLogger(__name__)


def split_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """简单文本分块：按段落+句子边界切分，滑动窗口保持上下文连贯"""
    if not text or not text.strip():
        return []

    # 按段落拆分
    paragraphs = re.split(r'\n\s*\n', text.strip())
    chunks = []
    current = ''

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) + 1 <= chunk_size:
            current = (current + '\n' + para).strip()
        else:
            if current:
                chunks.append(current)
                # 滑动窗口重叠
                if len(current) > overlap:
                    current = current[-overlap:] + '\n' + para
                else:
                    current = para
            else:
                # 单段超过 chunk_size，按句子切分
                sentences = re.split(r'(?<=[。！？.!?])\s*', para)
                for sent in sentences:
                    sent = sent.strip()
                    if not sent:
                        continue
                    if len(current) + len(sent) + 1 <= chunk_size:
                        current = (current + '\n' + sent).strip()
                    else:
                        if current:
                            chunks.append(current)
                        current = sent
    if current:
        chunks.append(current)

    return chunks if chunks else [text.strip()]


def estimate_tokens(text: str) -> int:
    """粗略估算 token 数量（中文 ~1.5 字符/token，英文 ~4 字符/token）"""
    chinese = len(re.findall(r'[一-鿿]', text))
    other = len(text) - chinese
    return int(chinese / 1.5 + other / 4)


ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'doc', 'docx', 'txt', 'md'}


class DocumentService:

    def save_upload(self, file: FileStorage, course_id: int, uploader_id: int, title: str = None) -> Courseware:
        """保存上传的文件，创建 courseware 记录"""
        # 校验文件类型
        ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(f'不支持的文件格式: .{ext}')

        # 生成安全文件名
        filename = secure_filename(file.filename)
        # 按课程分目录
        upload_dir = self._get_upload_dir(course_id)
        os.makedirs(upload_dir, exist_ok=True)

        # 保存文件
        file_path = os.path.join(upload_dir, filename)
        file.save(file_path)

        # 创建数据库记录
        courseware = Courseware(
            course_id=course_id,
            title=title or os.path.splitext(file.filename)[0],
            file_type=ext,
            file_path=file_path,
            file_size=os.path.getsize(file_path),
            uploader_id=uploader_id,
            status='uploading',
        )
        db.session.add(courseware)
        db.session.commit()

        return courseware

    def get_file_path(self, courseware_id: int) -> str:
        """获取课件文件路径"""
        cw = db.session.get(Courseware, courseware_id)
        if not cw:
            raise NotFoundError('课件不存在')
        return cw.file_path

    def extract_text(self, file_path: str, file_type: str) -> tuple:
        """
        从文件中提取纯文本。
        返回 (full_text, page_count)
        """
        if file_type == 'txt' or file_type == 'md':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(), 1

        elif file_type == 'pdf':
            return self._extract_pdf(file_path)

        elif file_type in ('ppt', 'pptx'):
            return self._extract_pptx(file_path)

        elif file_type in ('doc', 'docx'):
            return self._extract_docx(file_path)

        else:
            raise ValueError(f'不支持的文件类型: {file_type}')

    def process_document(self, courseware_id: int, chunk_size: int = 500, chunk_overlap: int = 50):
        """
        处理课件：提取文本 → 分块 → 写入 MySQL document_chunks 表。
        """
        cw = db.session.get(Courseware, courseware_id)
        if not cw:
            raise NotFoundError('课件不存在')

        try:
            cw.status = 'processing'
            db.session.commit()

            # 1. 提取文本
            full_text, page_count = self.extract_text(cw.file_path, cw.file_type)
            cw.page_count = page_count

            # 2. 文本分块
            chunks = split_text(full_text, chunk_size, chunk_overlap)

            # 3. 写入 document_chunks 表
            for idx, chunk_text in enumerate(chunks):
                chunk = DocumentChunk(
                    courseware_id=courseware_id,
                    chunk_index=idx,
                    content=chunk_text,
                    token_count=estimate_tokens(chunk_text),
                )
                db.session.add(chunk)

            cw.chunk_count = len(chunks)
            cw.status = 'completed'
            db.session.commit()

        except Exception as e:
            cw.status = 'failed'
            db.session.commit()
            logger.error(f'文档处理失败 courseware_id={courseware_id}: {e}')
            raise

    def _extract_pdf(self, file_path: str) -> tuple:
        """PDF 文本提取"""
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return '\n\n'.join(pages), len(pages)

    def _extract_pptx(self, file_path: str) -> tuple:
        """PPT/PPTX 文本提取"""
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append('\n'.join(texts))
        return '\n\n'.join(slides), len(slides)

    def _extract_docx(self, file_path: str) -> tuple:
        """DOCX 文本提取"""
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paragraphs), len(paragraphs)

    def _get_upload_dir(self, course_id: int) -> str:
        """获取课程上传目录"""
        from flask import current_app
        base = current_app.config['UPLOAD_FOLDER']
        return os.path.join(base, f'course_{course_id}')


document_service = DocumentService()
